#!/usr/bin/env python3
"""
Raíces Desarrolladores — Presentación para Inversores
10 slides / python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Brand colors ────────────────────────────────────────────────────
INK    = RGBColor(0x1A, 0x1A, 0x18)
LINO   = RGBColor(0xF0, 0xED, 0xE6)
BLANCO = RGBColor(0xF8, 0xF7, 0xF4)
SUELO  = RGBColor(0xE8, 0xE3, 0xDA)
TIERRA = RGBColor(0x85, 0x60, 0x40)
ARENA  = RGBColor(0xC4, 0xA9, 0x7D)
LIQUEN = RGBColor(0x7E, 0x8E, 0x5F)
MUSGO  = RGBColor(0x53, 0x62, 0x45)
CEIBO  = RGBColor(0xB8, 0x5C, 0x42)
NIEBLA = RGBColor(0x8F, 0x90, 0x89)
GREEN_DIM = RGBColor(0x36, 0x3D, 0x2A)

SERIF = 'Georgia'
SANS  = 'Helvetica Neue'

W = Inches(13.333)
H = Inches(7.5)

LOGO = '/Users/dariomintzer/Documents/Raices-Web/public/assets/images/raices_logo_transparente.png'
OUT  = '/Users/dariomintzer/Documents/Raices-Web/Raices-Inversores.pptx'


# ── Helpers ─────────────────────────────────────────────────────────

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, fill=None, line=None, lw=Pt(0.75)):
    s = slide.shapes.add_shape(1, l, t, w, h)
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line
        s.line.width = lw
    else:
        s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h,
        font=SANS, size=12, bold=False, italic=False,
        color=INK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def txt2(slide, lines, l, t, w, h, default_align=PP_ALIGN.LEFT):
    """Multi-paragraph textbox. lines = list of dicts with keys:
       text, font, size, bold, italic, color, align, space_before"""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = item.get('align', default_align)
        if 'space_before' in item:
            p.space_before = Pt(item['space_before'])
        run = p.add_run()
        run.text = item['text']
        run.font.name = item.get('font', SANS)
        run.font.size = Pt(item.get('size', 12))
        run.font.bold = item.get('bold', False)
        run.font.italic = item.get('italic', False)
        run.font.color.rgb = item.get('color', INK)
    return tb

def logo_img(slide, l, t, h):
    if os.path.exists(LOGO):
        try:
            return slide.shapes.add_picture(LOGO, l, t, height=h)
        except Exception:
            pass

def hline(slide, l, t, w, color, h=Pt(0.75)):
    return rect(slide, l, t, w, h, fill=color)

def vline(slide, l, t, h, color, w=Pt(0.75)):
    return rect(slide, l, t, w, h, fill=color)

def slide_num(slide, n, color=NIEBLA):
    txt(slide, f'{n:02d} / 10',
        W - Inches(0.8), H - Inches(0.42), Inches(0.65), Inches(0.28),
        font=SANS, size=7.5, color=color, align=PP_ALIGN.RIGHT)

def brand_footer(slide, color=NIEBLA):
    txt(slide, 'RAÍCES DESARROLLADORES',
        Inches(0.45), H - Inches(0.42), Inches(4), Inches(0.28),
        font=SANS, size=7.5, color=color)

def small_label(slide, text, l, t, w, color=MUSGO):
    txt(slide, text, l, t, w, Inches(0.32),
        font=SANS, size=7.5, bold=True, color=color)


# ── SLIDE 1 · PORTADA ───────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

s = blank(prs)
bg(s, INK)

vline(s, Inches(0), Inches(0), H, MUSGO, w=Inches(0.07))
rect(s, Inches(9.2), Inches(0), Inches(4.133), Inches(0.055), fill=ARENA)

logo_img(s, Inches(0.45), Inches(0.44), Inches(1.05))

txt(s, 'PRESENTACIÓN PARA INVERSORES',
    Inches(0.45), Inches(1.82), Inches(7), Inches(0.3),
    font=SANS, size=7.5, bold=True, color=ARENA)
hline(s, Inches(0.45), Inches(2.16), Inches(2.6), ARENA)

txt2(s, [
    {'text': 'Naturaleza', 'font': SERIF, 'size': 74, 'color': BLANCO},
    {'text': 'urbana.', 'font': SERIF, 'size': 74, 'color': LIQUEN},
], Inches(0.45), Inches(2.25), Inches(9), Inches(3.0))

txt(s, 'Edificios con arraigo, verde integrado e identidad propia.',
    Inches(0.45), Inches(5.28), Inches(6.6), Inches(0.55),
    font=SANS, size=13.5, color=NIEBLA)

hline(s, Inches(0.45), Inches(6.58), Inches(4.8), RGBColor(0x36, 0x3D, 0x2A))
txt(s, 'RAÍCES DESARROLLADORES · BUENOS AIRES · 2025',
    Inches(0.45), Inches(6.72), Inches(7), Inches(0.32),
    font=SANS, size=7.5, color=MUSGO)

for i in range(6):
    hline(s, Inches(10.4), Inches(2.6 + i * 0.32), Inches(2.5),
          RGBColor(0x2C, 0x33, 0x20), h=Pt(0.75))


# ── SLIDE 2 · LA OPORTUNIDAD ────────────────────────────────────────
s = blank(prs)
bg(s, LINO)

rect(s, Inches(0), Inches(0), W, Inches(0.08), fill=MUSGO)
vline(s, Inches(0), Inches(0), H, MUSGO, w=Inches(0.07))

small_label(s, 'LA OPORTUNIDAD', Inches(0.55), Inches(0.32), Inches(5), color=MUSGO)

txt(s, 'El mercado construye igual.\nRaíces, no.',
    Inches(0.55), Inches(0.75), Inches(8.5), Inches(1.9),
    font=SERIF, size=40, color=INK)

hline(s, Inches(0.55), Inches(2.65), Inches(12.2), SUELO)

cols_data = [
    ('01', 'Exceso de oferta genérica',
     'El mercado porteño tiene sobreoferta de proyectos intercambiables. La competencia es por precio; no existe diferenciación de marca.'),
    ('02', 'Verde como decoración',
     'La vegetación se agrega al final del proyecto como recurso de venta. No hay sistema real: sin tensores, sin riego, sin mantenimiento previsto.'),
    ('03', 'Desarrolladoras sin narrativa',
     'Las desarrolladoras no construyen marca entre proyectos. Cada obra es una entidad aislada, sin continuidad ni reconocimiento acumulado.'),
]
for i, (num, title, body) in enumerate(cols_data):
    x = Inches(0.55 + i * 4.26)
    rect(s, x, Inches(2.88), Inches(3.95), Inches(3.75),
         fill=BLANCO, line=SUELO, lw=Pt(0.5))
    rect(s, x, Inches(2.88), Inches(3.95), Inches(0.055),
         fill=[MUSGO, ARENA, CEIBO][i])
    txt(s, num, x + Inches(0.25), Inches(3.08), Inches(1), Inches(0.65),
        font=SERIF, size=34, color=ARENA)
    txt(s, title, x + Inches(0.25), Inches(3.8), Inches(3.45), Inches(0.5),
        font=SANS, size=12.5, bold=True, color=INK)
    txt(s, body, x + Inches(0.25), Inches(4.38), Inches(3.45), Inches(1.65),
        font=SANS, size=10.5, color=TIERRA)

brand_footer(s)
slide_num(s, 2)


# ── SLIDE 3 · FILOSOFÍA ─────────────────────────────────────────────
s = blank(prs)
bg(s, BLANCO)

rect(s, Inches(0), Inches(0), Inches(5.35), H, fill=INK)
vline(s, Inches(0), Inches(0), H, MUSGO, w=Inches(0.07))

logo_img(s, Inches(0.38), Inches(0.44), Inches(0.95))

small_label(s, 'FILOSOFÍA', Inches(0.38), Inches(1.72), Inches(4.5), color=ARENA)
hline(s, Inches(0.38), Inches(2.06), Inches(2.8), ARENA)

txt(s, 'La ciudad también\npuede echar raíces.',
    Inches(0.38), Inches(2.2), Inches(4.7), Inches(2.3),
    font=SERIF, size=36, color=BLANCO)

txt(s, 'Cada proyecto es una oportunidad de dejar\nalgo vivo en el entorno. No edificamos\npara el mercado genérico: edificamos\ncon criterio, identidad y futuro.',
    Inches(0.38), Inches(4.6), Inches(4.65), Inches(1.5),
    font=SANS, size=11, color=NIEBLA)

principles = [
    ('Naturaleza con criterio',
     'El verde entra en la decisión de fachada, balcones, terraza, riego y mantenimiento. Desde el primer boceto, no al final.'),
    ('Identidad acumulada',
     'Cada proyecto nombrado como un árbol autóctono. La marca crece con cada obra; el reconocimiento es colectivo y duradero.'),
    ('Rentabilidad con propósito',
     'Si invertir en verde reduce el margen, sigue siendo el camino correcto. El objetivo es una desarrolladora con criterio propio y futuro.'),
]
for i, (title, body) in enumerate(principles):
    y = Inches(0.95 + i * 2.1)
    if i > 0:
        hline(s, Inches(5.65), y - Inches(0.22), Inches(7.3), SUELO)
    vline(s, Inches(5.52), y, Inches(1.45), LIQUEN, w=Inches(0.04))
    txt(s, title, Inches(5.72), y, Inches(7.1), Inches(0.48),
        font=SANS, size=13.5, bold=True, color=INK)
    txt(s, body, Inches(5.72), y + Inches(0.54), Inches(7.1), Inches(1.1),
        font=SANS, size=11, color=TIERRA)

slide_num(s, 3, ARENA)


# ── SLIDE 4 · SISTEMA DE MARCA ──────────────────────────────────────
s = blank(prs)
bg(s, SUELO)

rect(s, Inches(0), Inches(0), W, Inches(0.08), fill=MUSGO)
vline(s, Inches(0), Inches(0), H, MUSGO, w=Inches(0.07))

small_label(s, 'IDENTIDAD VISUAL', Inches(0.55), Inches(0.32), Inches(6), color=MUSGO)

txt(s, 'Un sistema visual\nbautizado como un bosque.',
    Inches(0.55), Inches(0.72), Inches(7.5), Inches(1.62),
    font=SERIF, size=34, color=INK)

swatches = [
    (INK,    'Tinta',  '#1A1A18'),
    (MUSGO,  'Musgo',  '#536245'),
    (LIQUEN, 'Liquen', '#7E8E5F'),
    (ARENA,  'Arena',  '#C4A97D'),
    (LINO,   'Lino',   '#F0EDE6'),
    (CEIBO,  'Ceibo',  '#B85C42'),
]
for i, (color, name, hex_) in enumerate(swatches):
    x = Inches(0.55 + i * 2.06)
    swatch_w = Inches(1.82)
    rect(s, x, Inches(2.58), swatch_w, Inches(0.88), fill=color,
         line=ARENA if color == LINO else None, lw=Pt(0.5))
    txt(s, name,  x, Inches(3.55), swatch_w, Inches(0.28),
        font=SANS, size=8.5, bold=True, color=INK)
    txt(s, hex_,  x, Inches(3.83), swatch_w, Inches(0.28),
        font=SANS, size=7.5, color=NIEBLA)

hline(s, Inches(0.55), Inches(4.28), Inches(12.22), ARENA)

small_label(s, 'SISTEMA DE NAMING', Inches(0.55), Inches(4.52), Inches(6), color=MUSGO)

txt(s, 'Cada proyecto = un árbol autóctono + la calle.',
    Inches(0.55), Inches(4.95), Inches(8.5), Inches(0.48),
    font=SERIF, size=20, color=INK)

naming_cards = [
    ('Ceibo Vidal',         'P.001 · Activo',   True),
    ('Jacarandá',           'P.002 · Reserva',  False),
    ('Ombú',                'P.003 · Reserva',  False),
    ('Tipa',                'P.004 · Reserva',  False),
    ('Lapacho · Algarrobo…','Cartera futura',   False),
]
for i, (name, sub, active) in enumerate(naming_cards):
    x = Inches(0.55 + i * 2.54)
    fill_c = BLANCO if active else RGBColor(0xDC, 0xD7, 0xCE)
    rect(s, x, Inches(5.6), Inches(2.35), Inches(0.72),
         fill=fill_c, line=ARENA, lw=Pt(0.5))
    rect(s, x, Inches(5.6), Inches(2.35), Inches(0.055),
         fill=CEIBO if active else NIEBLA)
    txt(s, name, x + Inches(0.14), Inches(5.72), Inches(2.08), Inches(0.35),
        font=SANS, size=10.5, bold=active, color=INK if active else NIEBLA)
    txt(s, sub,  x + Inches(0.14), Inches(6.06), Inches(2.08), Inches(0.24),
        font=SANS, size=7.5, color=CEIBO if active else NIEBLA)

brand_footer(s)
slide_num(s, 4)


# ── SLIDE 5 · CEIBO VIDAL ───────────────────────────────────────────
s = blank(prs)
bg(s, BLANCO)

rect(s, Inches(0), Inches(0), W, Inches(2.72), fill=INK)
vline(s, Inches(0), Inches(0), H, CEIBO, w=Inches(0.07))

small_label(s, 'PROYECTO 001 · EN DESARROLLO', Inches(0.45), Inches(0.35), Inches(8), color=CEIBO)
hline(s, Inches(0.45), Inches(0.72), Inches(2.4), ARENA)

txt(s, 'Ceibo Vidal',
    Inches(0.45), Inches(0.85), Inches(9), Inches(1.05),
    font=SERIF, size=56, color=BLANCO)

txt(s, 'Vidal 3849 · Palermo · Ciudad de Buenos Aires',
    Inches(0.45), Inches(1.92), Inches(7.5), Inches(0.52),
    font=SANS, size=12, color=ARENA)

stats = [
    ('PB + 4 pisos',   'Escala barrial',             MUSGO),
    ('3 tipologías',   'por planta tipo',             LIQUEN),
    ('39 – 79 m²',    'Monoambiente a 3 ambientes',  ARENA),
    ('100% con balcón','Expansión vegetal exterior',  CEIBO),
]
for i, (val, desc, accent) in enumerate(stats):
    x = Inches(0.45 + i * 3.22)
    rect(s, x, Inches(3.05), Inches(3.0), Inches(2.05),
         fill=LINO, line=SUELO, lw=Pt(0.5))
    rect(s, x, Inches(3.05), Inches(3.0), Inches(0.055), fill=accent)
    txt(s, val,  x + Inches(0.2), Inches(3.26), Inches(2.72), Inches(0.7),
        font=SERIF, size=24, color=INK)
    txt(s, desc, x + Inches(0.2), Inches(4.0), Inches(2.72), Inches(0.5),
        font=SANS, size=10, color=NIEBLA)

hline(s, Inches(0.45), Inches(5.28), Inches(12.4), SUELO)

txt(s, 'Estado: En desarrollo · Inicio de obra previsto 2025.',
    Inches(0.45), Inches(5.45), Inches(8.5), Inches(0.38),
    font=SANS, size=11, color=TIERRA)

txt(s, '"El primer edificio de una serie. El punto de partida de la marca."',
    Inches(0.45), Inches(5.98), Inches(11.5), Inches(0.45),
    font=SERIF, size=15.5, italic=True, color=INK)

brand_footer(s)
slide_num(s, 5)


# ── SLIDE 6 · VERDE EDILICIO ────────────────────────────────────────
s = blank(prs)
bg(s, MUSGO)

rect(s, Inches(6.88), Inches(0), Inches(6.453), H, fill=BLANCO)
vline(s, Inches(6.88), Inches(0), H, LIQUEN, w=Inches(0.07))

small_label(s, 'VERDE EDILICIO', Inches(0.45), Inches(0.44), Inches(5), color=ARENA)

txt(s, 'Un sistema vivo,\nno una maceta suelta.',
    Inches(0.45), Inches(1.05), Inches(5.9), Inches(2.1),
    font=SERIF, size=38, color=BLANCO)

txt(s, 'La vegetación se diseña con tensores, drenajes\ny riego automatizado desde el inicio del proyecto,\nno como agregado de último momento.',
    Inches(0.45), Inches(3.25), Inches(5.9), Inches(1.2),
    font=SANS, size=12, color=ARENA)

hline(s, Inches(0.45), Inches(4.6), Inches(3.0), ARENA)

txt(s, 'Diferencial real y sostenible\nfrente a la competencia.',
    Inches(0.45), Inches(4.78), Inches(5.9), Inches(0.75),
    font=SANS, size=12, italic=True, color=RGBColor(0xC4, 0xA9, 0x7D))

features = [
    ('Terrazas verdes',
     'Espacios comunes en la cubierta con jardineras, agua y sombra natural. El remate del edificio es habitable, no técnico.'),
    ('Balcones vivos',
     'Plantas trepadoras con tensores verticales desde el proyecto. Fachada vegetal a escala edilicia, no macetas aisladas.'),
    ('Sistema integrado',
     'Riego automatizado, drenajes y mantenimiento previstos desde el diseño. Sostenible para propietarios y administración.'),
]
for i, (title, body) in enumerate(features):
    y = Inches(0.95 + i * 2.1)
    if i > 0:
        hline(s, Inches(7.22), y - Inches(0.2), Inches(5.8), SUELO)
    vline(s, Inches(7.08), y + Inches(0.04), Inches(1.42), LIQUEN, w=Inches(0.04))
    txt(s, title, Inches(7.28), y, Inches(5.7), Inches(0.46),
        font=SANS, size=13.5, bold=True, color=INK)
    txt(s, body,  Inches(7.28), y + Inches(0.54), Inches(5.7), Inches(1.08),
        font=SANS, size=11, color=TIERRA)

slide_num(s, 6, ARENA)


# ── SLIDE 7 · TIPOLOGÍAS ────────────────────────────────────────────
s = blank(prs)
bg(s, BLANCO)

rect(s, Inches(0), Inches(0), W, Inches(0.08), fill=MUSGO)
vline(s, Inches(0), Inches(0), H, MUSGO, w=Inches(0.07))

small_label(s, 'TIPOLOGÍAS · CEIBO VIDAL', Inches(0.55), Inches(0.32), Inches(8), color=MUSGO)

txt(s, 'Tres tipologías por planta,\ntodas con balcón.',
    Inches(0.55), Inches(0.72), Inches(8.5), Inches(1.5),
    font=SERIF, size=34, color=INK)

headers = ['Tipología', 'Sup. cubierta', 'Balcón', 'Total', 'Por planta']
col_w = [Inches(3.0), Inches(2.3), Inches(1.6), Inches(1.9), Inches(2.5)]
col_x = [Inches(0.55)]
for w in col_w[:-1]:
    col_x.append(col_x[-1] + w)

rect(s, Inches(0.55), Inches(2.58), Inches(11.8), Inches(0.52), fill=INK)
for j, (h_txt, cx) in enumerate(zip(headers, col_x)):
    txt(s, h_txt, cx + Inches(0.14), Inches(2.67), col_w[j], Inches(0.32),
        font=SANS, size=9, bold=True, color=BLANCO)

rows = [
    ('Monoambiente (Dep. B)', '35,79 m²', '3,27 m²', '39,06 m²', '1 por piso'),
    ('2 ambientes (Dep. C)',  '48,39 m²', '7,79 m²', '56,18 m²', '1 por piso'),
    ('3 ambientes (Dep. A)',  '68,67 m²', '11,08 m²','79,75 m²', '1 por piso'),
]
for ri, row_data in enumerate(rows):
    y = Inches(3.1 + ri * 0.72)
    row_bg = BLANCO if ri % 2 == 0 else LINO
    rect(s, Inches(0.55), y, Inches(11.8), Inches(0.66),
         fill=row_bg, line=SUELO, lw=Pt(0.5))
    for j, (cell, cx) in enumerate(zip(row_data, col_x)):
        txt(s, cell, cx + Inches(0.14), y + Inches(0.15), col_w[j], Inches(0.38),
            font=SANS, size=11, bold=(j == 0), color=INK)

hline(s, Inches(0.55), Inches(5.3), Inches(11.8), SUELO)

txt(s, 'Planta baja + 4 pisos · Núcleo central de circulación · Todas las unidades tienen expansión exterior.',
    Inches(0.55), Inches(5.48), Inches(12), Inches(0.38),
    font=SANS, size=10.5, color=NIEBLA)

txt(s, '"Escala barrial que maximiza calidad espacial por planta."',
    Inches(0.55), Inches(6.08), Inches(10.5), Inches(0.45),
    font=SERIF, size=15.5, italic=True, color=TIERRA)

brand_footer(s)
slide_num(s, 7)


# ── SLIDE 8 · PROYECCIÓN DE MARCA ───────────────────────────────────
s = blank(prs)
bg(s, INK)

vline(s, Inches(0), Inches(0), H, LIQUEN, w=Inches(0.07))

small_label(s, 'VISIÓN A LARGO PLAZO', Inches(0.45), Inches(0.44), Inches(6), color=LIQUEN)

txt(s, 'Una marca que crece\nproyecto a proyecto.',
    Inches(0.45), Inches(0.9), Inches(8.5), Inches(1.85),
    font=SERIF, size=42, color=BLANCO)

txt(s, 'El modelo es acumulativo: cada edificio refuerza el anterior.\nEl naming por árboles construye reconocimiento colectivo y valor de marca.',
    Inches(0.45), Inches(2.85), Inches(8.8), Inches(0.75),
    font=SANS, size=12, color=NIEBLA)

# Timeline horizontal
hline(s, Inches(0.45), Inches(3.88), Inches(12.45), GREEN_DIM)

projects = [
    ('2025',  'Ceibo Vidal',         'P.001 · Palermo · Activo', True,  CEIBO),
    ('2026',  'Jacarandá',           'P.002 · En reserva',       False, LIQUEN),
    ('2027',  'Ombú',                'P.003 · En reserva',       False, ARENA),
    ('2028+', 'Tipa · Lapacho · …', 'Cartera en expansión',     False, NIEBLA),
]
for i, (year, name, sub, active, color) in enumerate(projects):
    x = Inches(0.45 + i * 3.28)
    dot_size = Inches(0.24) if active else Inches(0.16)
    dot_offset = Inches(0) if active else Inches(0.04)
    rect(s, x + dot_offset, Inches(3.75) + dot_offset, dot_size, dot_size,
         fill=color)
    txt(s, year, x, Inches(4.18), Inches(3.1), Inches(0.32),
        font=SANS, size=8.5, bold=True, color=color)
    txt(s, name, x, Inches(4.55), Inches(3.1), Inches(0.65),
        font=SERIF, size=22, color=BLANCO if active else GREEN_DIM)
    txt(s, sub,  x, Inches(5.28), Inches(3.1), Inches(0.35),
        font=SANS, size=9.5, color=ARENA if active else NIEBLA)

hline(s, Inches(0.45), Inches(6.38), Inches(5.5), ARENA)
txt(s, 'La rentabilidad por proyecto se suma a la construcción de una marca con valor independiente.',
    Inches(0.45), Inches(6.54), Inches(11.5), Inches(0.42),
    font=SANS, size=10.5, italic=True, color=NIEBLA)

slide_num(s, 8, ARENA)


# ── SLIDE 9 · POR QUÉ INVERTIR ──────────────────────────────────────
s = blank(prs)
bg(s, LINO)

rect(s, Inches(0), Inches(0), W, Inches(0.08), fill=CEIBO)
vline(s, Inches(0), Inches(0), H, CEIBO, w=Inches(0.07))

small_label(s, 'POR QUÉ INVERTIR EN RAÍCES', Inches(0.55), Inches(0.32), Inches(8), color=CEIBO)

txt(s, 'Cuatro razones para\nser parte desde el inicio.',
    Inches(0.55), Inches(0.72), Inches(8.5), Inches(1.58),
    font=SERIF, size=36, color=INK)

reasons = [
    ('Diferenciación real',
     'No competimos por precio. Competimos por identidad, calidad de diseño y verde integrado. Eso protege márgenes y permite pricing propio.',
     CEIBO),
    ('Marca con valor acumulado',
     'Cada proyecto construye la siguiente venta. El nombre Raíces vale más con cada edificio entregado. La marca es el activo.',
     MUSGO),
    ('Escala controlada, riesgo bajo',
     'PB + 4 pisos reduce la exposición financiera sin resignar calidad. Proyectos ejecutables, plazos reales, sin especulación.',
     LIQUEN),
    ('Ventana de entrada temprana',
     'El primer proyecto define las condiciones de la sociedad. Entrar ahora es participar de la construcción de la marca desde cero.',
     ARENA),
]
for i, (title, body, accent) in enumerate(reasons):
    col, row = i % 2, i // 2
    x = Inches(0.55 + col * 6.42)
    y = Inches(2.62 + row * 2.18)
    rect(s, x, y, Inches(6.08), Inches(1.95),
         fill=BLANCO, line=SUELO, lw=Pt(0.5))
    rect(s, x, y, Inches(6.08), Inches(0.055), fill=accent)
    txt(s, title, x + Inches(0.22), y + Inches(0.22), Inches(5.64), Inches(0.45),
        font=SANS, size=13, bold=True, color=INK)
    txt(s, body,  x + Inches(0.22), y + Inches(0.72), Inches(5.64), Inches(0.96),
        font=SANS, size=10.5, color=TIERRA)

brand_footer(s)
slide_num(s, 9)


# ── SLIDE 10 · CIERRE ───────────────────────────────────────────────
s = blank(prs)
bg(s, INK)

vline(s, Inches(6.62), Inches(0), H, MUSGO, w=Inches(0.07))
vline(s, Inches(0), Inches(0), H, MUSGO, w=Inches(0.07))

# Left
small_label(s, 'RAÍCES DESARROLLADORES', Inches(0.45), Inches(1.62), Inches(5.6), color=ARENA)
hline(s, Inches(0.45), Inches(1.98), Inches(3.8), ARENA)

txt(s, 'Construyamos\njuntos.',
    Inches(0.45), Inches(2.18), Inches(5.8), Inches(2.3),
    font=SERIF, size=54, color=BLANCO)

logo_img(s, Inches(0.45), Inches(4.75), Inches(1.28))

txt(s, 'Esta presentación es confidencial y fue preparada\nexclusivamente para el destinatario.',
    Inches(0.45), Inches(6.52), Inches(5.8), Inches(0.55),
    font=SANS, size=8, italic=True, color=NIEBLA)

# Right — contact
small_label(s, 'CONTACTO', Inches(7.0), Inches(1.62), Inches(5.8), color=ARENA)
hline(s, Inches(7.0), Inches(1.98), Inches(5.92), ARENA)

contacts = [
    ('Email',        'contacto@raicesdesarrolladores.com.ar'),
    ('WhatsApp',     '+54 9 11 · disponible en la web'),
    ('Web',          'raicesdesarrolladores.com.ar'),
    ('Sede',         'Buenos Aires, Argentina'),
]
for i, (lbl, val) in enumerate(contacts):
    y = Inches(2.22 + i * 0.9)
    txt(s, lbl.upper(), Inches(7.0), y, Inches(2.4), Inches(0.3),
        font=SANS, size=7.5, bold=True, color=NIEBLA)
    txt(s, val, Inches(7.0), y + Inches(0.3), Inches(5.92), Inches(0.42),
        font=SANS, size=12, color=BLANCO)

hline(s, Inches(7.0), Inches(5.95), Inches(5.92), GREEN_DIM)

txt2(s, [
    {'text': 'RAÍCES', 'font': SANS, 'size': 28, 'bold': False,
     'color': RGBColor(0x2C, 0x33, 0x20)},
], Inches(7.0), Inches(6.2), Inches(5.9), Inches(0.8))

slide_num(s, 10, ARENA)


# ── SAVE ────────────────────────────────────────────────────────────
prs.save(OUT)
print(f'✓ Guardado → {OUT}')
